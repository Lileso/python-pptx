# encoding: utf-8

"""
Test suite for pptx.action module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.action import ActionSetting, Hyperlink
from pptx.enum.action import PP_ACTION

from .unitutil.cxml import element
from .unitutil.mock import class_mock, instance_mock, property_mock


class DescribeActionSetting(object):

    def it_knows_its_action_type(self, action_fixture):
        action_setting, expected_action = action_fixture
        action = action_setting.action
        assert action is expected_action

    def it_provides_access_to_its_hyperlink(self, hyperlink_fixture):
        action_setting, hyperlink_, Hyperlink_ = hyperlink_fixture[:3]
        xPr, parent = hyperlink_fixture[3:]
        hyperlink = action_setting.hyperlink
        Hyperlink_.assert_called_once_with(xPr, parent, False)
        assert hyperlink is hyperlink_

    def it_can_find_its_slide_jump_target(self, target_slide_fixture):
        action_setting, expected_value = target_slide_fixture
        target_slide = action_setting.target_slide
        assert target_slide == expected_value

    def it_raises_on_no_next_prev_slide(self, tgt_sld_raise_fixture):
        action_setting = tgt_sld_raise_fixture
        with pytest.raises(ValueError):
            action_setting.target_slide

    def it_knows_its_slide_index_to_help(self, _slide_index_fixture):
        action_setting, slides, slide, expected_value = _slide_index_fixture
        slide_index = action_setting._slide_index
        slides.index.assert_called_once_with(slide)
        assert slide_index == expected_value

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=[
        ('p:cNvPr',              None,
         PP_ACTION.NONE),
        ('p:cNvPr/a:hlinkClick', None,
         PP_ACTION.HYPERLINK),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=firstslide',
         PP_ACTION.FIRST_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=lastslide',
         PP_ACTION.LAST_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=nextslide',
         PP_ACTION.NEXT_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=previousslide',
         PP_ACTION.PREVIOUS_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkshowjump?jump=endshow',
         PP_ACTION.END_SHOW),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinksldjump',
         PP_ACTION.NAMED_SLIDE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkfile',
         PP_ACTION.OPEN_FILE),
        ('p:cNvPr/a:hlinkClick', 'ppaction://hlinkpres',
         PP_ACTION.PLAY),
        ('p:cNvPr/a:hlinkClick', 'ppaction://customshow',
         PP_ACTION.NAMED_SLIDE_SHOW),
        ('p:cNvPr/a:hlinkClick', 'ppaction://ole',
         PP_ACTION.OLE_VERB),
        ('p:cNvPr/a:hlinkClick', 'ppaction://macro',
         PP_ACTION.RUN_MACRO),
        ('p:cNvPr/a:hlinkClick', 'ppaction://program',
         PP_ACTION.RUN_PROGRAM),
        ('p:cNvPr/a:hlinkClick',
         'ppaction://hlinkshowjump?jump=lastslideviewed',
         PP_ACTION.LAST_SLIDE_VIEWED),
    ])
    def action_fixture(self, request):
        cNvPr_cxml, action_text, expected_action = request.param
        cNvPr = element(cNvPr_cxml)
        if action_text is not None:
            cNvPr.hlinkClick.action = action_text
        action_setting = ActionSetting(cNvPr, None)
        return action_setting, expected_action

    @pytest.fixture
    def hyperlink_fixture(self, Hyperlink_, hyperlink_):
        xPr, parent = 'xPr', 'parent'
        action_setting = ActionSetting(xPr, parent)
        return action_setting, hyperlink_, Hyperlink_, xPr, parent

    @pytest.fixture
    def _slide_index_fixture(self, request, part_prop_):
        action_setting = ActionSetting(None, None)
        slide = part_prop_.return_value
        slides = slide.package.presentation.slides
        expected_value = 123
        slides.index.return_value = expected_value
        return action_setting, slides, slide, expected_value

    @pytest.fixture(params=[
        (PP_ACTION.NONE,                None),
        (PP_ACTION.HYPERLINK,           None),
        (PP_ACTION.FIRST_SLIDE,         0),
        (PP_ACTION.LAST_SLIDE,          5),
        (PP_ACTION.NEXT_SLIDE,          3),
        (PP_ACTION.PREVIOUS_SLIDE,      1),
        (PP_ACTION.END_SHOW,            None),
        (PP_ACTION.NAMED_SLIDE,         4),
        (PP_ACTION.OPEN_FILE,           None),
        (PP_ACTION.PLAY,                None),
        (PP_ACTION.NAMED_SLIDE_SHOW,    None),
        (PP_ACTION.OLE_VERB,            None),
        (PP_ACTION.RUN_MACRO,           None),
        (PP_ACTION.RUN_PROGRAM,         None),
        (PP_ACTION.LAST_SLIDE_VIEWED,   None),
    ])
    def target_slide_fixture(self, request, action_prop_, _slide_index_prop_,
                             part_prop_):
        action_type, expected_value = request.param
        cNvPr = element('p:cNvPr/a:hlinkClick{r:id=rId6}')
        action_setting = ActionSetting(cNvPr, None)
        action_prop_.return_value = action_type
        _slide_index_prop_.return_value = 2
        # this becomes the return value of ActionSetting._slides
        part_prop_.return_value.package.presentation.slides = [
            0, 1, 2, 3, 4, 5
        ]
        part_prop_.return_value.rels.related_parts = {'rId6': 4}
        return action_setting, expected_value

    @pytest.fixture(params=[
        (PP_ACTION.NEXT_SLIDE,     2),
        (PP_ACTION.PREVIOUS_SLIDE, 0),
    ])
    def tgt_sld_raise_fixture(self, request, action_prop_, part_prop_,
                              _slide_index_prop_):
        action_type, slide_idx = request.param
        action_setting = ActionSetting(None, None)
        action_prop_.return_value = action_type
        # this becomes the return value of ActionSetting._slides
        part_prop_.return_value.package.presentation.slides = [0, 1, 2]
        _slide_index_prop_.return_value = slide_idx
        return action_setting

    # fixture components ---------------------------------------------

    @pytest.fixture
    def action_prop_(self, request):
        return property_mock(request, ActionSetting, 'action')

    @pytest.fixture
    def Hyperlink_(self, request, hyperlink_):
        return class_mock(
            request, 'pptx.action.Hyperlink', return_value=hyperlink_
        )

    @pytest.fixture
    def hyperlink_(self, request):
        return instance_mock(request, Hyperlink)

    @pytest.fixture
    def part_prop_(self, request):
        return property_mock(request, ActionSetting, 'part')

    @pytest.fixture
    def _slide_index_prop_(self, request):
        return property_mock(request, ActionSetting, '_slide_index')